import logging
from collections import deque
from io import BytesIO
from typing import Final, override

from pptx import Presentation

from clever_faq.application.common.ports.document.file_processor import FileProcessor
from clever_faq.infrastructure.errors.file import CantReadFileError

logger: Final[logging.Logger] = logging.getLogger(__name__)


class PptxFileProcessor(FileProcessor):
    @override
    def extract_text(self, file_bytes: bytes) -> str:
        try:
            presentation: Presentation = Presentation(BytesIO(file_bytes))
            text_content: deque[str] = deque()

            # Обработка каждого слайда
            for slide in presentation.slides:
                # Обработка заголовков
                if slide.shapes.title:
                    text_content.append(slide.shapes.title.text)

                # Обработка всех фигур на слайде
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    # Обработка каждого абзаца в фигуре
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text: deque[str] = deque()
                        # Обработка каждого фрагмента текста (run)
                        for run in paragraph.runs:
                            paragraph_text.append(run.text)

                        if paragraph_text:
                            text_content.append("".join(paragraph_text))

            return "\n".join(text_content)

        except Exception as e:
            logger.exception("Error processing PPTX file")
            msg = "Invalid PPTX file format"
            raise CantReadFileError(msg) from e
